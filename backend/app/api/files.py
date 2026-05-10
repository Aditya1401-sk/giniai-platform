from fastapi import APIRouter, UploadFile, File
import PyPDF2
import pytesseract
from PIL import Image
import io
import os
from rag_service.rag_core import rag_engine

router = APIRouter()

@router.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    filename = file.filename
    content_type = file.content_type
    extracted_text = ""

    # 1. Handle PDF
    if filename.endswith(".pdf"):
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(await file.read()))
        for page in pdf_reader.pages:
            extracted_text += page.extract_text() + "\n"

    # 2. Handle Images (JPG/PNG) via OCR
    elif content_type.startswith("image/"):
        try:
            image = Image.open(io.BytesIO(await file.read()))
            extracted_text = pytesseract.image_to_string(image)
        except Exception as e:
            return {"error": "OCR system (Tesseract) is not installed on this server. Please install it to index images."}

    # 3. Handle Word Documents (.docx)
    elif filename.endswith(".docx"):
        import docx
        from lxml import etree
        try:
            doc = docx.Document(io.BytesIO(await file.read()))
            text_parts = []

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Extract all table cells
            for table in doc.tables:
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        text_parts.append(" | ".join(row_texts))

            # Extract headers and footers
            for section in doc.sections:
                if section.header:
                    for para in section.header.paragraphs:
                        if para.text.strip():
                            text_parts.append(para.text)
                if section.footer:
                    for para in section.footer.paragraphs:
                        if para.text.strip():
                            text_parts.append(para.text)

            extracted_text = "\n".join(text_parts)

            # Fallback: use raw XML text extraction if nothing found
            if not extracted_text.strip():
                xml_content = doc.element.xml
                root = etree.fromstring(xml_content.encode())
                extracted_text = " ".join(root.itertext())

        except Exception as e:
            return {"error": f"Failed to parse Word Document: {str(e)}"}


    # 4. Handle PowerPoint (.pptx)
    elif filename.endswith(".pptx"):
        from pptx import Presentation
        try:
            prs = Presentation(io.BytesIO(await file.read()))
            text_parts = []
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = para.text.strip()
                            if line:
                                slide_texts.append(line)
                if slide_texts:
                    text_parts.append(f"[Slide {i+1}]\n" + "\n".join(slide_texts))
            extracted_text = "\n\n".join(text_parts)
        except Exception as e:
            return {"error": f"Failed to parse PowerPoint: {str(e)}"}

    # 5. Fallback: Try reading as plain text (CSV, Markdown, Code, etc.)
    else:
        try:
            extracted_text = (await file.read()).decode("utf-8")
        except UnicodeDecodeError:
            return {"error": f"Unsupported binary file format ({filename}). Please upload text-based files, PDFs, Word Docs, or Images."}

    if extracted_text.strip():
        # Add to RAG knowledge base
        content_with_meta = f"Document Filename: {filename}\n\n{extracted_text}"
        rag_engine.add_document(content_with_meta, {"filename": filename})
        return {
            "message": f"Successfully indexed {filename}",
            "text_preview": extracted_text[:100] + "..."
        }
    
    return {"error": "No text could be extracted from the file"}
